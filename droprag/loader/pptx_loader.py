"""DropRAG PPTX 加载器 - python-pptx"""

import os
from typing import Optional, List
from droprag.loader import LoaderBase, LoadedDocument, _get_folder_info, _get_file_mtime
from droprag.logging import get_logger

log = get_logger(__name__)


class PptxLoader(LoaderBase):
    extensions = [".pptx", ".ppt"]

    def load(self, filepath: str, base_path: str) -> Optional[LoadedDocument]:
        try:
            from pptx import Presentation
        except ImportError:
            log.warning("python-pptx 未安装，跳过 PPTX 文件: pip install droprag[office]")
            return None

        try:
            prs = Presentation(filepath)
        except Exception as e:
            log.debug(f"PPTX 加载失败: {filepath} ({e})")
            return None

        slides = []
        heading = ""

        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- 幻灯片 {i} ---"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                            if not heading and i == 1:
                                heading = text

                # 提取表格
                if shape.has_table:
                    table_text = self._table_to_text(shape.table)
                    if table_text.strip():
                        parts.append("[表格]\n" + table_text)

            # 提取备注
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[备注] {notes}")

            slides.append("\n".join(parts))

        content = "\n\n".join(slides)
        if not content.strip():
            return None

        folder, subfolder = _get_folder_info(filepath, base_path)
        return LoadedDocument(
            content=content,
            source=filepath,
            filename=os.path.basename(filepath),
            file_type="pptx",
            category="",
            folder=folder,
            subfolder=subfolder,
            file_size=os.path.getsize(filepath),
            file_mtime=_get_file_mtime(filepath),
            heading=heading,
        )

    @staticmethod
    def _table_to_text(table) -> str:
        """将表格转换为文本格式"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
